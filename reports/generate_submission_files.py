from pathlib import Path
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches

project_root = Path(__file__).resolve().parents[1]
reports_dir = project_root / 'reports'

# Read report markdown
report_md = reports_dir / 'Final_Report.md'
report_pdf = reports_dir / 'Final_Report_generated.pdf'
text = report_md.read_text(encoding='utf-8')

# Generate PDF from markdown text
c = canvas.Canvas(str(report_pdf), pagesize=letter)
width, height = letter
margin = 50
line_height = 14
x = margin
y = height - margin
for line in text.splitlines():
    if y < margin:
        c.showPage()
        y = height - margin
    c.drawString(x, y, line[:95])
    y -= line_height
c.save()

# Build PowerPoint from outline and content
presentation = Presentation()
outline = reports_dir / 'Presentation_Outline.md'
content = reports_dir / 'Presentation_Content.md'

for slide_info in content.read_text(encoding='utf-8').split('\n\n'):
    if slide_info.strip():
        lines = slide_info.strip().splitlines()
        title = lines[0].replace('# ', '').strip()
        body = '\n'.join(line.strip('- ').strip() for line in lines[1:])
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

pptx_path = reports_dir / 'Presentation_generated.pptx'
presentation.save(str(pptx_path))

print(f'Generated PDF: {report_pdf.name}')
print(f'Generated PPTX: {pptx_path.name}')
